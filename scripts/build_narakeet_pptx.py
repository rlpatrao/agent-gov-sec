"""
scripts/build_narakeet_pptx.py

Parses docs/video-demo-script.md and builds a Narakeet-ready PPTX.

Each ### section becomes one slide:
  - Slide title     = section heading (e.g. "2.1 — Generating a Trace (10:00 – 11:00)")
  - Slide body      = [ON SCREEN] description + code/ASCII blocks (placeholder text
                      that you replace with your screenshot in PowerPoint)
  - Speaker notes   = all *"..."* narration lines joined into a paragraph

Usage:
    uv pip install python-pptx
    uv run python scripts/build_narakeet_pptx.py

Output:
    docs/galaxy-showcase-narakeet.pptx
    (open in PowerPoint, paste screenshots over placeholder text on each slide,
     then upload to Narakeet to generate the narrated video)
"""
from __future__ import annotations

import re
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt
except ImportError:
    print("ERROR: python-pptx not installed.")
    print("Run:  uv pip install python-pptx")
    sys.exit(1)

SCRIPT_PATH = Path(__file__).parent.parent / "docs" / "video-demo-script.md"
OUTPUT_PATH = Path(__file__).parent.parent / "docs" / "galaxy-showcase-narakeet.pptx"

# Slide dimensions: widescreen 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Colours
BG_DARK   = RGBColor(0x1E, 0x1E, 0x2E)   # dark background
TITLE_CLR = RGBColor(0xCB, 0xD6, 0xF7)   # light lavender
BODY_CLR  = RGBColor(0xA6, 0xE3, 0xA1)   # green text (placeholder hint)
NOTE_CLR  = RGBColor(0x00, 0x00, 0x00)


def _add_textbox(slide, left, top, width, height, text, font_size, color, bold=False, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    return txBox


def _set_bg(slide, color: RGBColor):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _parse_sections(md: str) -> list[dict]:
    """Split on ### headings; return list of {title, body, narration}."""
    parts = re.split(r"^### ", md, flags=re.MULTILINE)
    slides = []
    for part in parts[1:]:
        lines = part.split("\n")
        title = lines[0].strip()
        rest = "\n".join(lines[1:])

        # Collect narration: *"..."* patterns (single or multi-line quoted)
        narration_lines = re.findall(r'\*"([^"]+)"\*', rest, re.DOTALL)
        narration = " ".join(line.strip().replace("\n", " ") for line in narration_lines)

        # Collect [ON SCREEN] body: everything that isn't pure narration
        # Strip narration italic quotes, strip markdown noise, keep code blocks
        body_lines = []
        in_fence = False
        for line in rest.split("\n"):
            if line.strip().startswith("```"):
                in_fence = not in_fence
                if not in_fence:
                    continue
                body_lines.append("")
                continue
            if in_fence:
                body_lines.append(line)
                continue
            # Skip pure narration lines
            if re.match(r'^\s*\*"', line):
                continue
            # Skip fenced-block language hints already handled above
            # Keep [ON SCREEN], [TERMINAL], [NARRATE] tags and plain text
            cleaned = re.sub(r"`\[([^\]]+)\]`", r"[\1]", line)
            cleaned = re.sub(r"\*\*([^*]+)\*\*", r"\1", cleaned)   # bold
            cleaned = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", cleaned)  # links
            if cleaned.strip():
                body_lines.append(cleaned)

        body = "\n".join(body_lines).strip()
        # Truncate very long bodies — presenter replaces with screenshot anyway
        if len(body) > 1200:
            body = body[:1200] + "\n…"

        slides.append({"title": title, "body": body, "narration": narration})

    return slides


def build_pptx(script_path: Path, output_path: Path) -> None:
    md = script_path.read_text(encoding="utf-8")
    sections = _parse_sections(md)

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    blank_layout = prs.slide_layouts[6]  # completely blank

    for sec in sections:
        slide = prs.slides.add_slide(blank_layout)
        _set_bg(slide, BG_DARK)

        # Title bar (top strip)
        _add_textbox(
            slide,
            left=Inches(0.3), top=Inches(0.15),
            width=Inches(12.7), height=Inches(0.6),
            text=sec["title"],
            font_size=18, color=TITLE_CLR, bold=True,
        )

        # Body (placeholder — presenter pastes screenshot here)
        body_text = sec["body"] if sec["body"] else "[SCREENSHOT — see screenshot-manifest.md]"
        _add_textbox(
            slide,
            left=Inches(0.3), top=Inches(0.85),
            width=Inches(12.7), height=Inches(6.3),
            text=body_text,
            font_size=10, color=BODY_CLR,
        )

        # Speaker notes
        if sec["narration"]:
            notes_slide = slide.notes_slide
            tf = notes_slide.notes_text_frame
            tf.text = sec["narration"]
            for para in tf.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(14)

    prs.save(output_path)
    print(f"Saved {len(sections)} slides → {output_path}")


if __name__ == "__main__":
    if not SCRIPT_PATH.exists():
        print(f"ERROR: script not found at {SCRIPT_PATH}")
        sys.exit(1)
    build_pptx(SCRIPT_PATH, OUTPUT_PATH)
